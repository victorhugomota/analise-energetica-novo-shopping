#!/usr/bin/env python3
"""
Sincronizador Completo e Oficial: Site/Firebase -> PowerPoint (PPTX)
Atualiza 100% dos dados:
- Parâmetros Gerais (Firebase config/parametros)
- Tags e Potências dos Equipamentos (Firebase equipamentos e perfil_uso)
- Slide 6: Tabela de Parâmetros Baseline
- Slide 7: Gráficos de Sensibilidade (Tarifa e Duty Cycle Escrava) + Textos Explicativos
- Slide 10: Tabela Comparativa Consolidada + Gráficos de Barras (Consumo Mensal e Gasto Anual)
- Todos os slides: Textos de Economia, Payback, ROI, VPL e ESG CO2
"""
import os
import sys
import json
import urllib.request
import pptx
from pptx.chart.data import CategoryChartData
import win32com.client
import pythoncom

if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

PPTX_PATH = os.path.abspath(r"C:\Users\Eletrica\Desktop\Estudo de Custo Novo Shopping\Analise Final\Estudo_Viabilidade_Retrofit_HVAC_Novo_Shopping.pptx")
EXCEL_PATH = os.path.abspath(r"C:\Users\Eletrica\Desktop\Estudo de Custo Novo Shopping\Final\Planilha_Retrofit.xlsx")

def fetch_data_from_firebase():
    dados_completos_url = "https://firestore.googleapis.com/v1/projects/analiseenergeticanovoshopping/databases/(default)/documents/config/dados_completos"
    params_url = "https://firestore.googleapis.com/v1/projects/analiseenergeticanovoshopping/databases/(default)/documents/config/parametros"
    equip_url = "https://firestore.googleapis.com/v1/projects/analiseenergeticanovoshopping/databases/(default)/documents/equipamentos?pageSize=100"
    perfil_url = "https://firestore.googleapis.com/v1/projects/analiseenergeticanovoshopping/databases/(default)/documents/perfil_uso?pageSize=100"

    params = {
        'setpoint': 22.0, 'histerese': 1.0, 'erroDuplo': 2.5, 'tempoEscrava': 1.0,
        'tempoRevezamento': 12.0, 'tarifa': 0.75, 'capex': 301000.0,
        'dcAntes': 0.70, 'dcMestre': 0.70, 'dcEscrava': 0.50
    }
    custom_equip = []
    custom_perfil = []

    # 1. Tentar ler do documento consolidado de alta performance (dados_completos)
    try:
        req = urllib.request.Request(dados_completos_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode('utf-8'))
            fields = data.get('fields', {})
            # Parse parametros
            if 'parametros' in fields:
                p_map = fields['parametros'].get('mapValue', {}).get('fields', {})
                for k in params:
                    if k in p_map:
                        params[k] = float(p_map[k].get('doubleValue') or p_map[k].get('integerValue') or params[k])
            # Parse equipamentos
            if 'equipamentos' in fields:
                arr = fields['equipamentos'].get('arrayValue', {}).get('values', [])
                for i, v in enumerate(arr):
                    ef = v.get('mapValue', {}).get('fields', {})
                    custom_equip.append({
                        'idx': i,
                        'tag': ef.get('tag', {}).get('stringValue') or f'UE-{i+1:02d}',
                        'potUC': float(ef.get('potUC', {}).get('doubleValue') or ef.get('potUC', {}).get('integerValue') or 10.0),
                        'qtdUC': float(ef.get('qtdUC', {}).get('doubleValue') or ef.get('qtdUC', {}).get('integerValue') or 2.0),
                        'potUE': float(ef.get('potUE', {}).get('doubleValue') or ef.get('potUE', {}).get('integerValue') or 1.5),
                        'qtdUE': float(ef.get('qtdUE', {}).get('doubleValue') or ef.get('qtdUE', {}).get('integerValue') or 1.0),
                    })
            # Parse perfil
            if 'perfil_uso' in fields:
                arr = fields['perfil_uso'].get('arrayValue', {}).get('values', [])
                for i, v in enumerate(arr):
                    uf = v.get('mapValue', {}).get('fields', {})
                    custom_perfil.append({
                        'idx': i,
                        'tag': uf.get('tag', {}).get('stringValue') or f'UE-{i+1:02d}',
                        'hSegSex': float(uf.get('hSegSex', {}).get('doubleValue') or uf.get('hSegSex', {}).get('integerValue') or 12.0),
                        'hSab': float(uf.get('hSab', {}).get('doubleValue') or uf.get('hSab', {}).get('integerValue') or 12.0),
                        'hDom': float(uf.get('hDom', {}).get('doubleValue') or uf.get('hDom', {}).get('integerValue') or 12.0),
                        'setupAntes': float(uf.get('setupAntes', {}).get('doubleValue') or uf.get('setupAntes', {}).get('integerValue') or 1.5),
                        'setupDepois': float(uf.get('setupDepois', {}).get('doubleValue') or uf.get('setupDepois', {}).get('integerValue') or 1.5),
                    })
            if custom_equip and len(custom_equip) == 38:
                print(f"[OK] 38 Equipamentos e Parametros carregados com sucesso do documento consolidado Firebase!")
                return params, custom_equip, custom_perfil
    except Exception as e:
        pass

    # 2. Fallback para documentos legados individuais se o consolidado ainda nao foi criado
    try:
        req = urllib.request.Request(params_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode('utf-8'))
            fields = data.get('fields', {})
            params.update({
                'setpoint': float(fields.get('setpoint', {}).get('doubleValue') or fields.get('setpoint', {}).get('integerValue') or 22.0),
                'histerese': float(fields.get('histerese', {}).get('doubleValue') or fields.get('histerese', {}).get('integerValue') or 1.0),
                'erroDuplo': float(fields.get('erroDuplo', {}).get('doubleValue') or fields.get('erroDuplo', {}).get('integerValue') or 2.5),
                'tempoEscrava': float(fields.get('tempoEscrava', {}).get('doubleValue') or fields.get('tempoEscrava', {}).get('integerValue') or 1.0),
                'tempoRevezamento': float(fields.get('tempoRevezamento', {}).get('doubleValue') or fields.get('tempoRevezamento', {}).get('integerValue') or 12.0),
                'tarifa': float(fields.get('tarifa', {}).get('doubleValue') or fields.get('tarifa', {}).get('integerValue') or 0.75),
                'capex': float(fields.get('capex', {}).get('doubleValue') or fields.get('capex', {}).get('integerValue') or 301000.0),
                'dcAntes': float(fields.get('dcAntes', {}).get('doubleValue') or fields.get('dcAntes', {}).get('integerValue') or 0.70),
                'dcMestre': float(fields.get('dcMestre', {}).get('doubleValue') or fields.get('dcMestre', {}).get('integerValue') or 0.70),
                'dcEscrava': float(fields.get('dcEscrava', {}).get('doubleValue') or fields.get('dcEscrava', {}).get('integerValue') or 0.50),
            })
            print("[OK] Parametros carregados do Firebase!")
    except Exception as e:
        print(f"[AVISO] Parametros locais utilizados ({e})")

    try:
        req = urllib.request.Request(equip_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode('utf-8'))
            docs = data.get('documents', [])
            for doc in docs:
                f = doc.get('fields', {})
                custom_equip.append({
                    'idx': int(f.get('idx', {}).get('integerValue') or 0),
                    'tag': f.get('tag', {}).get('stringValue') or '',
                    'potUC': float(f.get('potUC', {}).get('doubleValue') or f.get('potUC', {}).get('integerValue') or 10.0),
                    'qtdUC': float(f.get('qtdUC', {}).get('doubleValue') or f.get('qtdUC', {}).get('integerValue') or 2.0),
                    'potUE': float(f.get('potUE', {}).get('doubleValue') or f.get('potUE', {}).get('integerValue') or 1.5),
                    'qtdUE': float(f.get('qtdUE', {}).get('doubleValue') or f.get('qtdUE', {}).get('integerValue') or 1.0),
                })
            custom_equip.sort(key=lambda x: x['idx'])
            if custom_equip:
                print(f"[OK] {len(custom_equip)} Equipamentos (com TAGs editadas) carregados do Firebase!")
    except Exception as e:
        pass

    try:
        req = urllib.request.Request(perfil_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode('utf-8'))
            docs = data.get('documents', [])
            for doc in docs:
                f = doc.get('fields', {})
                custom_perfil.append({
                    'idx': int(f.get('idx', {}).get('integerValue') or 0),
                    'tag': f.get('tag', {}).get('stringValue') or '',
                    'hSegSex': float(f.get('hSegSex', {}).get('doubleValue') or f.get('hSegSex', {}).get('integerValue') or 12.0),
                    'hSab': float(f.get('hSab', {}).get('doubleValue') or f.get('hSab', {}).get('integerValue') or 12.0),
                    'hDom': float(f.get('hDom', {}).get('doubleValue') or f.get('hDom', {}).get('integerValue') or 12.0),
                    'setupAntes': float(f.get('setupAntes', {}).get('doubleValue') or f.get('setupAntes', {}).get('integerValue') or 1.5),
                    'setupDepois': float(f.get('setupDepois', {}).get('doubleValue') or f.get('setupDepois', {}).get('integerValue') or 1.5),
                })
            custom_perfil.sort(key=lambda x: x['idx'])
    except Exception as e:
        pass

    return params, custom_equip, custom_perfil

def calculate_system(params, custom_equip=None, custom_perfil=None, custom_dc_escrava=None, custom_horas=None, custom_tarifa=None):
    tarifa = custom_tarifa if custom_tarifa is not None else params['tarifa']
    dc_escrava = custom_dc_escrava if custom_dc_escrava is not None else params['dcEscrava']

    if custom_equip and len(custom_equip) == 38:
        equipamentos = custom_equip
    else:
        import openpyxl
        wb = openpyxl.load_workbook(EXCEL_PATH, data_only=True)
        ws_c = wb["Cadastro de Equipamentos"]
        equipamentos = []
        for r in range(4, 42):
            tag = ws_c[f"A{r}"].value
            if not tag or tag == "Total": continue
            potUC = float(ws_c[f"B{r}"].value or 0)
            qtdUC = float(ws_c[f"C{r}"].value or 0)
            potUE = float(ws_c[f"D{r}"].value or 0)
            qtdUE = float(ws_c[f"E{r}"].value or 0)
            equipamentos.append({
                'tag': tag, 'potUC': potUC, 'qtdUC': qtdUC, 'potUE': potUE, 'qtdUE': qtdUE,
                'potTotal': potUC*qtdUC + potUE*qtdUE,
                'corrente': (potUC*qtdUC + potUE*qtdUE)*1000 / (1.732 * 380 * 0.95)
            })

    if custom_perfil and len(custom_perfil) == 38:
        perfil = []
        for p in custom_perfil:
            h_seg = custom_horas if custom_horas is not None else p['hSegSex']
            h_sab = custom_horas if custom_horas is not None else p['hSab']
            h_dom = custom_horas if custom_horas is not None else p['hDom']
            perfil.append({
                'tag': p['tag'], 'hSegSex': h_seg, 'hSab': h_sab, 'hDom': h_dom,
                'setupAntes': p['setupAntes'], 'setupDepois': p['setupDepois']
            })
    else:
        import openpyxl
        wb = openpyxl.load_workbook(EXCEL_PATH, data_only=True)
        ws_u = wb["Perfil de Uso e Consumo"]
        perfil = []
        for r in range(4, 42):
            tag = ws_u[f"A{r}"].value
            if not tag or "Total" in str(tag): continue
            h_seg = custom_horas if custom_horas is not None else float(ws_u[f"B{r}"].value or 12)
            h_sab = custom_horas if custom_horas is not None else float(ws_u[f"C{r}"].value or 12)
            h_dom = custom_horas if custom_horas is not None else float(ws_u[f"D{r}"].value or 12)
            perfil.append({
                'tag': tag, 'hSegSex': h_seg, 'hSab': h_sab, 'hDom': h_dom,
                'setupAntes': float(ws_u[f"E{r}"].value or 1.5),
                'setupDepois': float(ws_u[f"F{r}"].value or 1.5),
            })

    def consumo_dia(eq, horas, setup, modo):
        consUE = eq['potUE'] * eq['qtdUE'] * horas
        if horas <= 0: return consUE
        if modo == 'antes':
            consUC_setup = eq['qtdUC'] * eq['potUC'] * min(horas, setup)
            consUC_manut = eq['qtdUC'] * eq['potUC'] * max(0, horas - setup) * params['dcAntes']
            return consUE + consUC_setup + consUC_manut
        else:
            consUC_setup = eq['qtdUC'] * eq['potUC'] * min(horas, setup)
            consUC_manut = (1 * eq['potUC'] * params['dcMestre'] + (eq['qtdUC'] - 1) * eq['potUC'] * dc_escrava) * max(0, horas - setup)
            return consUE + consUC_setup + consUC_manut

    consAntesMes = 0.0
    consDepoisMes = 0.0
    for eq, u in zip(equipamentos, perfil):
        cA = 4.33 * (5 * consumo_dia(eq, u['hSegSex'], u['setupAntes'], 'antes') +
                     1 * consumo_dia(eq, u['hSab'], u['setupAntes'], 'antes') +
                     1 * consumo_dia(eq, u['hDom'], u['setupAntes'], 'antes'))
        cD = 4.33 * (5 * consumo_dia(eq, u['hSegSex'], u['setupDepois'], 'depois') +
                     1 * consumo_dia(eq, u['hSab'], u['setupDepois'], 'depois') +
                     1 * consumo_dia(eq, u['hDom'], u['setupDepois'], 'depois'))
        consAntesMes += cA
        consDepoisMes += cD

    econMesKwh = consAntesMes - consDepoisMes
    consAntesAno = consAntesMes * 12
    consDepoisAno = consDepoisMes * 12
    econAnoKwh = econMesKwh * 12
    redPerc = (econMesKwh / consAntesMes) * 100 if consAntesMes > 0 else 0.0

    custoAntesMes = consAntesMes * tarifa
    custoDepoisMes = consDepoisMes * tarifa
    econMesR = econMesKwh * tarifa

    custoAntesAno = custoAntesMes * 12
    custoDepoisAno = custoDepoisMes * 12
    econAnoR = econMesR * 12

    paybackMeses = params['capex'] / econMesR if econMesR > 0 else 999.0
    roi5 = ((econAnoR * 5 - params['capex']) / params['capex']) * 100 if params['capex'] > 0 else 0.0
    vpl = econAnoR * 3.604776 - params['capex']
    co2 = (econAnoKwh * 0.086) / 1000.0

    horasOperacao = perfil[0]['hSegSex'] if perfil else 12.0
    setupHoras = perfil[0]['setupAntes'] if perfil else 1.5

    return {
        'consAntesMes': consAntesMes, 'consDepoisMes': consDepoisMes, 'econMesKwh': econMesKwh,
        'consAntesAno': consAntesAno, 'consDepoisAno': consDepoisAno, 'econAnoKwh': econAnoKwh,
        'redPerc': redPerc,
        'custoAntesMes': custoAntesMes, 'custoDepoisMes': custoDepoisMes, 'econMesR': econMesR,
        'custoAntesAno': custoAntesAno, 'custoDepoisAno': custoDepoisAno, 'econAnoR': econAnoR,
        'paybackMeses': paybackMeses, 'roi5': roi5, 'vpl': vpl, 'co2': co2,
        'capex': params['capex'], 'tarifa': tarifa,
        'dcAntes': params['dcAntes'], 'dcMestre': params['dcMestre'], 'dcEscrava': dc_escrava,
        'setpoint': params['setpoint'], 'histerese': params['histerese'], 'tempoRevezamento': params['tempoRevezamento'],
        'horasOperacao': horasOperacao, 'setupHoras': setupHoras,
        'equipamentos': equipamentos
    }

def update_presentation_complete(calc, params, custom_equip, custom_perfil):
    if not os.path.exists(PPTX_PATH):
        print(f"[ERRO] Arquivo nao encontrado em {PPTX_PATH}")
        return False

    pythoncom.CoInitialize()
    try:
        ppt_app = win32com.client.GetActiveObject("PowerPoint.Application")
        for pres in list(ppt_app.Presentations):
            if os.path.normpath(pres.FullName) == os.path.normpath(PPTX_PATH):
                pres.Close()
    except Exception:
        pass

    prs = pptx.Presentation(PPTX_PATH)

    fmtCur = lambda v: f"R$ {v:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    fmtInt = lambda v: f"{round(v):,}".replace(",", ".")
    fmtDec = lambda v, d=1: f"{v:.{d}f}".replace(".", ",")

    s_tar_065 = calculate_system(params, custom_equip, custom_perfil, custom_tarifa=0.65)['econAnoR'] / 1000.0
    s_tar_075 = calculate_system(params, custom_equip, custom_perfil, custom_tarifa=0.75)['econAnoR'] / 1000.0
    s_tar_085 = calculate_system(params, custom_equip, custom_perfil, custom_tarifa=0.85)['econAnoR'] / 1000.0
    s_tar_100 = calculate_system(params, custom_equip, custom_perfil, custom_tarifa=1.00)['econAnoR'] / 1000.0

    s_dc_30 = calculate_system(params, custom_equip, custom_perfil, custom_dc_escrava=0.30)['econAnoR'] / 1000.0
    s_dc_40 = calculate_system(params, custom_equip, custom_perfil, custom_dc_escrava=0.40)['econAnoR'] / 1000.0
    s_dc_50 = calculate_system(params, custom_equip, custom_perfil, custom_dc_escrava=0.50)['econAnoR'] / 1000.0
    s_dc_60 = calculate_system(params, custom_equip, custom_perfil, custom_dc_escrava=0.60)['econAnoR'] / 1000.0

    calc_plus2h = calculate_system(params, custom_equip, custom_perfil, custom_horas=calc['horasOperacao'] + 2.0)
    econ_unit_mes = calc['econMesR'] / 38.0
    econ_unit_ano = calc['econAnoR'] / 38.0
    econ_10_ano = econ_unit_ano * 10.0

    # 1. Atualizar SLIDE 6 (Parâmetros Baseline)
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_table:
            t = shape.table
            for row in t.rows:
                p_name = row.cells[0].text.strip()
                if "Setpoint de Temperatura" in p_name:
                    row.cells[1].text = fmtDec(calc['setpoint'], 1)
                elif "Histerese de Controle" in p_name:
                    row.cells[1].text = fmtDec(calc['histerese'], 1)
                elif "Tempo de Revezamento" in p_name:
                    row.cells[1].text = fmtDec(calc['tempoRevezamento'], 1)
                elif "Tarifa de Energia" in p_name:
                    row.cells[1].text = fmtCur(calc['tarifa'])
                elif "Investimento CAPEX" in p_name:
                    row.cells[1].text = fmtCur(calc['capex'])
                elif "Duty Cycle Antes" in p_name:
                    row.cells[1].text = f"{calc['dcAntes']:.2f} ({round(calc['dcAntes']*100)}%)"
                elif "Duty Cycle Mestre" in p_name:
                    row.cells[1].text = f"{calc['dcMestre']:.2f} ({round(calc['dcMestre']*100)}%)"
                elif "Duty Cycle Escrava" in p_name:
                    row.cells[1].text = f"{calc['dcEscrava']:.2f} ({round(calc['dcEscrava']*100)}%)"
                elif "Horas de Operação" in p_name:
                    row.cells[1].text = fmtDec(calc['horasOperacao'], 1)
                elif "Tempo de Pull-Down" in p_name:
                    row.cells[1].text = fmtDec(calc['setupHoras'], 1)

    # 2. Atualizar SLIDE 7 (Sensibilidade, Gráficos e Textos)
    s7 = prs.slides[6]
    chart_count = 0
    for shape in s7.shapes:
        if shape.has_chart:
            chart_count += 1
            if chart_count == 1:
                cd = CategoryChartData()
                cd.categories = ['0,65', '0,75', '0,85', '1,00']
                cd.add_series('Economia Anual (R$ mil)', (round(s_tar_065, 1), round(s_tar_075, 1), round(s_tar_085, 1), round(s_tar_100, 1)))
                shape.chart.replace_data(cd)
            elif chart_count == 2:
                cd = CategoryChartData()
                cd.categories = ['30%', '40%', '50% (padrão)', '60%']
                cd.add_series('Economia Anual (R$ mil)', (round(s_dc_30, 1), round(s_dc_40, 1), round(s_dc_50, 1), round(s_dc_60, 1)))
                shape.chart.replace_data(cd)
        elif shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if "Horas de operação:" in p.text:
                    p.text = f"Horas de operação: +2h/dia (Black Friday/Natal) eleva a economia para {fmtCur(calc_plus2h['econAnoR'])}/ano (payback {fmtDec(calc_plus2h['paybackMeses'], 1)} meses).   CAPEX: fixado em {fmtCur(calc['capex'])} para as 38 unidades."
                elif "Escala linear por máquina" in p.text:
                    p.text = f"Escala linear por máquina retrofitada: cada climatizador gera {fmtCur(econ_unit_mes)}/mês ({fmtCur(econ_unit_ano)}/ano) de economia líquida individual — 10 máquinas = {fmtCur(econ_10_ano)}/ano; 38 máquinas (parque total) = {fmtCur(calc['econAnoR'])}/ano."

    # 3. Atualizar SLIDE 10 (Comparativo Consolidado e Gráficos)
    s10 = prs.slides[9]
    chart_count_10 = 0
    for shape in s10.shapes:
        if shape.has_table:
            t = shape.table
            for row in t.rows:
                p_name = row.cells[0].text.strip()
                if "Consumo Mensal de Energia" in p_name:
                    row.cells[1].text = f"{fmtInt(calc['consAntesMes'])} kWh"
                    row.cells[2].text = f"{fmtInt(calc['consDepoisMes'])} kWh"
                    row.cells[3].text = f"{fmtInt(calc['econMesKwh'])} kWh"
                    row.cells[4].text = f"-{fmtDec(calc['redPerc'], 1)}%"
                elif "Consumo Anual de Energia" in p_name:
                    row.cells[1].text = f"{fmtInt(calc['consAntesAno'])} kWh"
                    row.cells[2].text = f"{fmtInt(calc['consDepoisAno'])} kWh"
                    row.cells[3].text = f"{fmtInt(calc['econAnoKwh'])} kWh"
                    row.cells[4].text = f"-{fmtDec(calc['redPerc'], 1)}%"
                elif "Gasto Operacional Mensal" in p_name:
                    row.cells[1].text = fmtCur(calc['custoAntesMes'])
                    row.cells[2].text = fmtCur(calc['custoDepoisMes'])
                    row.cells[3].text = fmtCur(calc['econMesR'])
                    row.cells[4].text = f"-{fmtDec(calc['redPerc'], 1)}%"
                elif "Gasto Operacional Anual" in p_name:
                    row.cells[1].text = fmtCur(calc['custoAntesAno'])
                    row.cells[2].text = fmtCur(calc['custoDepoisAno'])
                    row.cells[3].text = fmtCur(calc['econAnoR'])
                    row.cells[4].text = f"-{fmtDec(calc['redPerc'], 1)}%"
        elif shape.has_chart:
            chart_count_10 += 1
            if chart_count_10 == 1:
                cd = CategoryChartData()
                cd.categories = ['Antes', 'Com CLP']
                cd.add_series('kWh/mês', (round(calc['consAntesMes']), round(calc['consDepoisMes'])))
                shape.chart.replace_data(cd)
            elif chart_count_10 == 2:
                cd = CategoryChartData()
                cd.categories = ['Antes', 'Com CLP']
                cd.add_series('R$/ano', (round(calc['custoAntesAno'], 2), round(calc['custoDepoisAno'], 2)))
                shape.chart.replace_data(cd)

    # 4. Atualizar textos em TODOS os slides (Global replacement)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = p.text
                    if "R$ 217.686" in txt or "R$ 217.686,42" in txt or "R$ 163.264" in txt or "R$ 163.264,82" in txt:
                        p.text = p.text.replace("R$ 217.686,42", fmtCur(calc['econAnoR'])).replace("R$ 217.686", fmtCur(calc['econAnoR']).split(',')[0]).replace("R$ 163.264,82", fmtCur(calc['econAnoR'])).replace("R$ 163.264", fmtCur(calc['econAnoR']).split(',')[0])
                    if "16,6 meses" in txt or "16,5 meses" in txt or "22,1 meses" in txt:
                        p.text = p.text.replace("16,6 meses", f"{fmtDec(calc['paybackMeses'], 1)} meses").replace("16,5 meses", f"{fmtDec(calc['paybackMeses'], 1)} meses").replace("22,1 meses", f"{fmtDec(calc['paybackMeses'], 1)} meses")
                    if "+261,6%" in txt or "+262,8%" in txt or "+171,2%" in txt or "+172,1%" in txt:
                        p.text = p.text.replace("+261,6%", f"+{fmtDec(calc['roi5'], 1)}%").replace("+262,8%", f"+{fmtDec(calc['roi5'], 1)}%").replace("+171,2%", f"+{fmtDec(calc['roi5'], 1)}%").replace("+172,1%", f"+{fmtDec(calc['roi5'], 1)}%")
                    if "25,0 tCO2/ano" in txt or "25,0 tCO" in txt or "18,7 tCO2/ano" in txt:
                        p.text = p.text.replace("25,0 tCO2/ano", f"{fmtDec(calc['co2'], 1)} tCO2/ano").replace("18,7 tCO2/ano", f"{fmtDec(calc['co2'], 1)} tCO2/ano")
                    if "290.249 kWh/ano" in txt or "217.686 kWh/ano" in txt:
                        p.text = p.text.replace("290.249 kWh/ano", f"{fmtInt(calc['econAnoKwh'])} kWh/ano").replace("217.686 kWh/ano", f"{fmtInt(calc['econAnoKwh'])} kWh/ano")
                    if "R$ 18.140,53" in txt or "R$ 18.140,54" in txt or "R$ 18.140" in txt or "R$ 13.605,40" in txt:
                        p.text = p.text.replace("R$ 18.140,53", fmtCur(calc['econMesR'])).replace("R$ 18.140,54", fmtCur(calc['econMesR'])).replace("R$ 18.140", fmtCur(calc['econMesR']).split(',')[0]).replace("R$ 13.605,40", fmtCur(calc['econMesR']))
                    if "R$ 483.710" in txt or "R$ 484.710" in txt or "R$ 287.525" in txt or "R$ 288.533" in txt:
                        p.text = p.text.replace("R$ 483.710", fmtCur(calc['vpl']).split(',')[0]).replace("R$ 484.710", fmtCur(calc['vpl']).split(',')[0]).replace("R$ 287.525", fmtCur(calc['vpl']).split(',')[0]).replace("R$ 288.533", fmtCur(calc['vpl']).split(',')[0])

    prs.save(PPTX_PATH)

    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    try:
        pres = ppt_app.Presentations.Open(PPTX_PATH, WithWindow=False)
        pres.Save()
        pres.Close()
        print("[SUCESSO] Apresentacao PowerPoint 100% sincronizada (Slides 6, 7, 10, tabelas e graficos)!")
        print(f"[ARQUIVO] {PPTX_PATH}")
        return True
    finally:
        ppt_app.Quit()
        pythoncom.CoUninitialize()

if __name__ == "__main__":
    params, custom_equip, custom_perfil = fetch_data_from_firebase()
    calc = calculate_system(params, custom_equip, custom_perfil)
    print("\n--- DADOS CONSOLIDADOS DO SISTEMA ---")
    print(f"• Tarifa: R$ {calc['tarifa']:.2f}/kWh | CAPEX: R$ {calc['capex']:,.2f}")
    print(f"• Duty Mestre: {calc['dcMestre']*100:.0f}% | Duty Escrava: {calc['dcEscrava']*100:.0f}%")
    print(f"• Horas de Operacao: {calc['horasOperacao']:.1f} h/dia | Setup: {calc['setupHoras']:.1f} h")
    print(f"• Economia Anual: R$ {calc['econAnoR']:,.2f} ({calc['econAnoKwh']:,.0f} kWh/ano)")
    print(f"• Economia Mensal: R$ {calc['econMesR']:,.2f} ({calc['econMesKwh']:,.0f} kWh/mês)")
    print(f"• Reducao Percentual: -{calc['redPerc']:.1f}%")
    print(f"• Payback Simples: {calc['paybackMeses']:.1f} meses | ROI (5a): +{calc['roi5']:.1f}% | VPL: R$ {calc['vpl']:,.2f}")
    print(f"• Descarbonizacao ESG: {calc['co2']:.1f} tCO2/ano")
    print("------------------------------------\n")
    update_presentation_complete(calc, params, custom_equip, custom_perfil)
